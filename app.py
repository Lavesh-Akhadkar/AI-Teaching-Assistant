import streamlit as st
from pptx import Presentation
from PyPDF2 import PdfReader
from llm import converse, add_document_to_store, delete_all_documents, clear_history
from youtube import video_to_transcript


if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" not in st.session_state:
    st.session_state.documents = False

if not st.session_state.documents:
    delete_all_documents()
    clear_history()
    st.session_state.documents = True


def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)

    for slide in presentation.slides:
        slide_text = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"

        add_document_to_store(slide_text)


def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            add_document_to_store(page_text)


def convert_audio_to_text(video_url):
    return video_to_transcript(video_url=video_url)


st.title("AI Teaching Assistant")

uploaded_files = st.file_uploader(
    "Upload a PowerPoint or PDF file", accept_multiple_files=True, type=["pptx", "pdf"]
)
uploaded_to_rag = st.button("Process Files")
video_url = st.text_input("Enter YouTube video URL")
process_video = st.button("Process Video")


if uploaded_files is not None and uploaded_to_rag:
    for uploaded_file in uploaded_files:
        if uploaded_file.name.endswith(".pptx"):
            with st.spinner("Extracting text from PowerPoint"):
                extract_text_from_pptx(uploaded_file)

        elif uploaded_file.name.endswith(".pdf"):
            with st.spinner("Extracting text from PDF"):
                extract_text_from_pdf(uploaded_file)

if video_url and process_video:
    with st.spinner("Processing video..."):
        text = convert_audio_to_text(video_url)
        add_document_to_store(text)


for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])


user_input = st.chat_input("Ask a question")
if user_input:
    with st.chat_message("user"):
        st.markdown(user_input)
    st.session_state.messages.append({"role": "user", "content": user_input})
    response = converse(user_input)
    with st.chat_message("ai"):
        st.markdown(response)
    st.session_state.messages.append({"role": "assistant", "content": response})
